"""Reproducible figure generation for the EJPS BBB manuscript.

All figures are generated from public data (data/descriptor_table.csv) and the
parameter file (data/parameters.csv).  Any illustrative scores are documented
explicitly in the Methods and figure legends.
"""
import csv
import math
import re
from pathlib import Path

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch, Patch, Rectangle
from matplotlib.colors import BoundaryNorm, ListedColormap
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


def load_csv(path):
    with open(path, newline='', encoding='utf-8') as f:
        return list(csv.DictReader(f))


def load_params(path):
    rows = load_csv(path)
    return {r['key']: r['value'] for r in rows}


def get_param(params, key, as_float=False):
    raw = str(params[key]).strip()
    if as_float:
        return float(raw)
    return raw


def _markdown_table(rows, widths=None):
    """Return a simple Markdown table from a list of rows."""
    def fmt(c):
        return str(c).replace('\n', ' ').replace('|', '\\|').strip()
    lines = []
    header = '| ' + ' | '.join(fmt(c) for c in rows[0]) + ' |'
    lines.append(header)
    lines.append('|' + '|'.join('-' * (len(fmt(c)) + 2) for c in rows[0]) + '|')
    for row in rows[1:]:
        cells = [fmt(row[i]) if i < len(row) else '' for i in range(len(rows[0]))]
        lines.append('| ' + ' | '.join(cells) + ' |')
    return '\n'.join(lines)


def _save_pptx(fig_paths, titles, captions, output_path, params):
    prs = Presentation()
    prs.slide_width = Inches(get_param(params, 'PPTX_WIDTH_IN', as_float=True))
    prs.slide_height = Inches(get_param(params, 'PPTX_HEIGHT_IN', as_float=True))

    for img_path, title, caption in zip(fig_paths, titles, captions):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.6))
        tf = title_box.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.font.size = Pt(20)
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT

        slide.shapes.add_picture(str(img_path), Inches(0.8), Inches(1.0), height=Inches(5.3))

        cap_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.8))
        tf2 = cap_box.text_frame
        tf2.text = caption
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.font.size = Pt(12)
        p2.alignment = PP_ALIGN.LEFT

    prs.save(str(output_path))


# ---------------------------------------------------------------------------
# Descriptor categorisation helpers
# ---------------------------------------------------------------------------
FAV_SYMBOLS = {'+': 'favorable', 'o': 'neutral', '-': 'weak', 'x': 'unfavorable'}


def _symbol_to_score(sym):
    return {'+': 2, 'o': 1, '-': 0, 'x': -1}.get(sym, 0)


def _a_d_symbol(ad, low, cutoff, high):
    if ad <= low:
        return '+'
    if ad <= cutoff:
        return 'o'
    if ad <= high:
        return '-'
    return 'x'


def _ccs_symbol(text):
    t = text.lower()
    if 'very large' in t or 'large' in t and 'medium' not in t:
        return 'x'
    if 'medium-to-large' in t or 'large' in t:
        return '-'
    if 'small-to-medium' in t or 'medium' in t:
        return 'o'
    if 'small' in t or 'minimal' in t:
        return '+'
    return 'o'


def _desolvation_symbol(text):
    t = text.lower()
    if 'high' in t:
        return 'x'
    if 'medium-to-high' in t:
        return '-'
    if 'medium' in t:
        return 'o'
    if 'low-to-medium' in t:
        return 'o'
    if 'low' in t:
        return '+'
    return 'o'


def _pgp_symbol(text):
    t = text.lower()
    if 'strong' in t:
        return 'x'
    if 'substrate' in t:
        return '-'
    if 'weak' in t:
        return 'o'
    if 'non/weak' in t:
        return 'o'
    if 'non' in t:
        return '+'
    return 'o'


def _dipole_symbol(text):
    t = text.lower()
    if 'high' in t:
        return '-'
    if 'medium' in t:
        return 'o'
    if 'low' in t:
        return '+'
    return 'o'


def _lumo_symbol(text):
    t = text.lower()
    # High LUMO energy is generally associated with lower reactivity; for BBB
    # this descriptor has only weak discrimination, so keep mapping mild.
    if 'high' in t:
        return '+'
    if 'medium' in t:
        return 'o'
    if 'low' in t:
        return '-'
    return 'o'


def _tpsa3d_symbol(text):
    t = text.lower()
    if 'very large' in t:
        return 'x'
    if 'large' in t:
        return '-'
    if 'medium' in t:
        return 'o'
    if 'small' in t:
        return '+'
    return 'o'


def _chameleon_symbol(text):
    t = text.lower()
    if 'medium' in t or 'low-to-medium' in t:
        return 'o'
    if 'low' in t or 'none' in t:
        return '-'
    if 'high' in t:
        return '+'
    return '-'


def _synergy_symbol(d):
    """Pattern descriptor: aromatic + amine/halogen combinations are common in BBB+ drugs."""
    text = (d.get('synergy', '') + ' ' + d.get('lateral_bilayer', '')).lower()
    if 'aromatic' in text and ('tertiary amine' in text or 'secondary amine' in text or 'amide' in text):
        return '+'
    if 'aromatic' in text:
        return 'o'
    return '-'


def _latpress_symbol(text):
    text = text.lower()
    if 'aromatic' in text and ('amine' in text or 'amide' in text):
        return '+'
    if 'aromatic' in text:
        return 'o'
    return '-'


# ---------------------------------------------------------------------------
# Figure 1: Discriminatory power ranking
# ---------------------------------------------------------------------------
def _power_to_score(power):
    p = power.lower()
    if 'unclear' in p or 'not applicable' in p:
        return 2.5, '#91bfdb', 'Unclear'
    if 'strong' in p and 'moderate' not in p:
        return 5.0, '#d73027', 'Strong'
    if 'moderate-to-strong' in p or 'moderate; high' in p or 'moderate to strong' in p:
        return 4.0, '#f46d43', 'Moderate-to-strong'
    if 'moderate' in p and ('limited' in p or 'weak' in p):
        return 2.5, '#91bfdb', 'Moderate/Weak'
    if 'moderate' in p:
        return 3.0, '#fc8d59', 'Moderate'
    if 'limited' in p or 'weak' in p or 'boundary' in p:
        return 2.0, '#91bfdb', 'Limited/Weak'
    return 2.5, '#91bfdb', 'Unclear'


def make_figure1(output_dir, table2_rows, params):
    # table2_rows is list of [Descriptor, Class, Definition, Key reference, Discriminatory power]
    descriptors = []
    scores = []
    colors = []
    labels = []
    for row in table2_rows[1:]:
        score, color, label = _power_to_score(row[4])
        descriptors.append(row[0])
        scores.append(score)
        colors.append(color)
        labels.append(label)

    # sort descending
    idx = np.argsort(scores)[::-1]
    descriptors = [descriptors[i] for i in idx]
    scores = [scores[i] for i in idx]
    colors = [colors[i] for i in idx]
    labels = [labels[i] for i in idx]

    fig, ax = plt.subplots(figsize=(8, 6))
    y = np.arange(len(descriptors))
    bars = ax.barh(y, scores, color=colors, edgecolor='black', linewidth=0.5)
    ax.set_yticks(y)
    ax.set_yticklabels(descriptors, fontsize=9)
    ax.invert_yaxis()
    ax.set_xlim(0, 6)
    ax.set_xlabel('Overall promise rating (schematic)', fontsize=10)
    ax.set_title('Figure 1. Discriminatory power ranking of unconventional and conventional descriptors', fontsize=11)

    # Add tier labels on bars
    for bar, score, label in zip(bars, scores, labels):
        ax.text(score + 0.1, bar.get_y() + bar.get_height()/2, label,
                va='center', ha='left', fontsize=8)

    # tier legend
    legend_elements = [
        Patch(facecolor='#d73027', edgecolor='black', label='Top tier (Strong)'),
        Patch(facecolor='#fc8d59', edgecolor='black', label='High tier (Moderate)'),
        Patch(facecolor='#91bfdb', edgecolor='black', label='Medium/Limited'),
    ]
    ax.legend(handles=legend_elements, loc='lower right', fontsize=8)
    fig.tight_layout()
    path = output_dir / 'figure1_discriminatory_power.png'
    fig.savefig(path, dpi=get_param(params, 'FIGURE_DPI', as_float=True), bbox_inches='tight')
    plt.close(fig)
    return path


# ---------------------------------------------------------------------------
# Figure 2: Drug x factor matrix
# ---------------------------------------------------------------------------
def make_figure2(output_dir, df, params):
    low = get_param(params, 'AD_LOW_A2', as_float=True)
    cutoff = get_param(params, 'AD_CUTOFF_A2', as_float=True)
    high = get_param(params, 'AD_HIGH_A2', as_float=True)

    factor_cols = [
        ('A_D', lambda d: _a_d_symbol(float(d['a_d']), low, cutoff, high)),
        ('CCS', lambda d: _ccs_symbol(d['ccs_trend'])),
        ('Cham.', lambda d: _chameleon_symbol(d['chameleon'])),
        ('Desolv.', lambda d: _desolvation_symbol(d['desolvation'])),
        ('Dipole', lambda d: _dipole_symbol(d['dipole'])),
        ('LUMO', lambda d: _lumo_symbol(d['lumo'])),
        ('Net Flux', lambda d: _pgp_symbol(d['p_gp'])),
        ('3D-PSA', lambda d: _tpsa3d_symbol(d['tpsa_3d'])),
        ('Substruct.', lambda d: _synergy_symbol(d)),
        ('Lat. Press.', lambda d: _latpress_symbol(d['lateral_bilayer'])),
    ]

    drugs = [d['drug'] for d in df]
    labels = [f[0] for f in factor_cols]
    matrix = np.zeros((len(df), len(factor_cols)), dtype=int)
    for i, d in enumerate(df):
        for j, (_, fn) in enumerate(factor_cols):
            matrix[i, j] = _symbol_to_score(fn(d))

    # discrete colormap aligned with symbol scores: x (-1), - (0), o (1), + (2)
    cmap = ListedColormap(['#a50026', '#f46d43', '#fee090', '#1a9850'])
    norm = BoundaryNorm(np.array([-1.5, -0.5, 0.5, 1.5, 2.5]), ncolors=4)

    fig, ax = plt.subplots(figsize=(10, 10))
    im = ax.imshow(matrix, cmap=cmap, norm=norm, aspect='auto')

    ax.set_xticks(np.arange(len(labels)))
    ax.set_yticks(np.arange(len(drugs)))
    ax.set_xticklabels(labels, rotation=45, ha='right', fontsize=9)
    ax.set_yticklabels(drugs, fontsize=9)

    # draw separator between BBB+ and BBB- groups
    n_pos = sum(1 for d in df if d['bbb_status'].startswith('+'))
    ax.axhline(y=n_pos - 0.5, color='black', linewidth=2)
    ax.text(len(labels) + 0.3, n_pos/2 - 0.5, 'BBB+', fontsize=10, color='#1a9850', fontweight='bold')
    ax.text(len(labels) + 0.3, n_pos + (len(drugs)-n_pos)/2 - 0.5, 'BBB-', fontsize=10, color='#a50026', fontweight='bold')

    # annotate with symbols
    symbol_map = {-1: 'x', 0: '−', 1: 'o', 2: '+'}
    for i in range(len(drugs)):
        for j in range(len(labels)):
            text = symbol_map.get(matrix[i, j], 'o')
            color = 'white' if matrix[i, j] in (-1, 2) else 'black'
            ax.text(j, i, text, ha='center', va='center', color=color, fontsize=10, fontweight='bold')

    ax.set_title('Figure 2. Drug × factor evaluation matrix (+: favorable, o: neutral, −: weak, x: unfavorable)', fontsize=11, pad=10)

    # colorbar
    cbar = fig.colorbar(im, ax=ax, fraction=0.03, pad=0.04)
    cbar.set_ticks([-1, 0, 1, 2])
    cbar.set_ticklabels(['Unfavorable', 'Weak', 'Neutral', 'Favorable'])
    cbar.ax.tick_params(labelsize=8)

    fig.tight_layout()
    path = output_dir / 'figure2_drug_factor_matrix.png'
    fig.savefig(path, dpi=get_param(params, 'FIGURE_DPI', as_float=True), bbox_inches='tight')
    plt.close(fig)
    return path


# ---------------------------------------------------------------------------
# Figure 3: A_D scatter
# ---------------------------------------------------------------------------
def make_figure3(output_dir, df, params):
    low = get_param(params, 'AD_LOW_A2', as_float=True)
    cutoff = get_param(params, 'AD_CUTOFF_A2', as_float=True)
    high = get_param(params, 'AD_HIGH_A2', as_float=True)

    # split into BBB+ and BBB- groups, order by A_D
    pos = [d for d in df if d['bbb_status'].startswith('+')]
    neg = [d for d in df if not d['bbb_status'].startswith('+')]

    fig, ax = plt.subplots(figsize=(10, 5))

    def plot_group(group, y, color, marker, edgecolor):
        np.random.seed(42)
        xs = [float(d['a_d']) for d in group]
        ys = [y + np.random.uniform(-0.08, 0.08) for _ in group]  # light jitter for visibility
        names = [d['drug'] for d in group]
        ax.scatter(xs, ys, c=color, marker=marker, s=120, edgecolors=edgecolor, linewidths=1, zorder=3)
        for x, yy, name in zip(xs, ys, names):
            ax.annotate(name, (x, yy), textcoords='offset points', xytext=(5, 5), fontsize=7, alpha=0.8)

    plot_group(pos, 1, '#1a9850', 'o', 'black')
    plot_group(neg, 0, '#d73027', 'x', 'black')

    ax.axvline(low, color='orange', linestyle='--', linewidth=1.2, label=f'$A_D \\approx {int(low)}$ Å² (easy crossing)')
    ax.axvline(cutoff, color='red', linestyle='--', linewidth=1.5, label=f'$A_D \\approx {int(cutoff)}$ Å² (cutoff)')
    ax.axvline(high, color='darkred', linestyle=':', linewidth=1.2, label=f'$A_D \\approx {int(high)}$ Å² (exclusion)')

    # shaded transition zone
    ax.axvspan(cutoff, high, alpha=0.1, color='orange', label='Transition zone')

    ax.set_yticks([0, 1])
    ax.set_yticklabels(['BBB- (impermeable)', 'BBB+ (permeable)'])
    ax.set_xlabel('Estimated membrane cross-sectional area $A_D$ (Å²)', fontsize=10)
    ax.set_ylabel('BBB permeability', fontsize=10)
    ax.set_title('Figure 3. BBB permeability as a function of estimated $A_D$', fontsize=11)
    ax.legend(loc='upper right', fontsize=8)
    ax.set_ylim(-0.5, 1.6)
    ax.set_xlim(0, max(float(d['a_d']) for d in df) * 1.1)
    fig.tight_layout()
    path = output_dir / 'figure3_ad_scatter.png'
    fig.savefig(path, dpi=get_param(params, 'FIGURE_DPI', as_float=True), bbox_inches='tight')
    plt.close(fig)
    return path


# ---------------------------------------------------------------------------
# Figure 4: Unified model schematic
# ---------------------------------------------------------------------------
def make_figure4(output_dir, params):
    fig, ax = plt.subplots(figsize=(9, 4.5))
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 5)
    ax.axis('off')

    def box(x, y, w, h, text, color, fontsize=10):
        rect = FancyBboxPatch((x, y), w, h, boxstyle='round,pad=0.05',
                              facecolor=color, edgecolor='black', linewidth=1.5)
        ax.add_patch(rect)
        ax.text(x + w/2, y + h/2, text, ha='center', va='center', fontsize=fontsize, wrap=True)

    box(0.3, 1.4, 2.4, 1.8, 'Step 1:\nDesolvation\n$P_{\\mathrm{desolv}}$', '#8dd3c7', fontsize=11)
    box(3.0, 1.4, 2.4, 1.8, 'Step 2:\nMembrane partition\n$P_{\\mathrm{partition}}$', '#fb8072', fontsize=11)
    box(5.7, 1.4, 2.4, 1.8, 'Step 3:\nNet flux\n$P_{\\mathrm{net\\ flux}}$', '#80b1d3', fontsize=11)
    box(8.4, 1.4, 1.5, 1.8, 'Brain\nparenchyma', '#bebada', fontsize=10)

    for x0, x1 in [(2.7, 3.0), (5.4, 5.7), (8.1, 8.4)]:
        ax.annotate('', xy=(x1, 2.3), xytext=(x0, 2.3),
                    arrowprops=dict(arrowstyle='->', lw=2, color='black'))

    # explanatory labels
    ax.text(1.5, 3.6, 'HBD / HBA / charge', ha='center', fontsize=8, style='italic')
    ax.text(4.2, 3.6, '$A_D$ / $\\pi_{bi}$ / logP', ha='center', fontsize=8, style='italic')
    ax.text(6.9, 3.6, 'P-gp / transporters', ha='center', fontsize=8, style='italic')

    ax.text(1.5, 0.8, 'desolvation cost', ha='center', fontsize=8)
    ax.text(4.2, 0.8, 'lateral bilayer pressure', ha='center', fontsize=8)
    ax.text(6.9, 0.8, 'efflux vs influx', ha='center', fontsize=8)

    ax.text(5, 4.4, r'$P_{BBB} \propto P_{\mathrm{desolv}} \times P_{\mathrm{partition}} \times P_{\mathrm{net\ flux}}$',
            ha='center', fontsize=12, fontweight='bold')

    fig.tight_layout()
    path = output_dir / 'figure4_unified_model.png'
    fig.savefig(path, dpi=get_param(params, 'FIGURE_DPI', as_float=True), bbox_inches='tight')
    plt.close(fig)
    return path


# ---------------------------------------------------------------------------
# Figure 5: Clinical paradoxes
# ---------------------------------------------------------------------------
def _clinical_case_scores(df, params):
    """Return a dict of illustrative component scores for selected clinical cases.

    The scores are relative, heuristic scores derived from the dataset:
    - P_desolv from the categorical desolvation column.
    - P_partition from the A_D relative term (exponential cutoff model).
    - P_net_flux from P-gp substrate status.
    """
    kb = get_param(params, 'BOLTZMANN_KB', as_float=True)
    T = get_param(params, 'TEMPERATURE_K', as_float=True)
    pi_mnm = get_param(params, 'BILAYER_PRESSURE_MNM', as_float=True)
    pi = pi_mnm / 1000.0
    A_ref = get_param(params, 'AD_REFERENCE_A2', as_float=True)
    kT = kb * T
    factor = pi / kT * 1e-20

    desolv_map = {
        'Low': 1.0, 'Low-to-medium': 0.8, 'Medium': 0.5,
        'Medium-to-high': 0.3, 'High': 0.15, 'Minimal': 1.0
    }

    pgp_map = {
        'Non': 1.0, 'Non/weak': 0.7, 'Weak': 0.5,
        'Substrate': 0.25, 'Strong': 0.05, 'Substrate (controversial)': 0.25
    }

    cases = ['Caffeine', 'Diazepam', 'Morphine', 'Heroin', 'Loperamide', 'Dopamine']
    rows = {}
    for d in df:
        if d['drug'] not in cases:
            continue
        ad = float(d['a_d'])
        p_partition = math.exp(-factor * (ad - A_ref))
        p_desolv = desolv_map.get(d['desolvation'].split(';')[0].strip(), 0.5)
        p_net = pgp_map.get(d['p_gp'].split(';')[0].strip(), 0.5)
        rows[d['drug']] = {
            'P_desolv': p_desolv,
            'P_partition': p_partition,
            'P_net_flux': p_net,
            'P_BBB': p_desolv * p_partition * p_net,
        }
    return rows


def make_figure5(output_dir, df, params):
    scores = _clinical_case_scores(df, params)
    cases = ['Caffeine', 'Diazepam', 'Morphine', 'Heroin', 'Loperamide', 'Dopamine']
    p_desolv = [scores[c]['P_desolv'] for c in cases]
    p_partition = [scores[c]['P_partition'] for c in cases]
    p_net = [scores[c]['P_net_flux'] for c in cases]
    p_bbb = [scores[c]['P_BBB'] for c in cases]

    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(10, 4.5))
    x = np.arange(len(cases))
    width = 0.25

    ax1.bar(x - width, p_desolv, width, label=r'$P_{\mathrm{desolv}}$', color='#8dd3c7', edgecolor='black')
    ax1.bar(x, p_partition, width, label=r'$P_{\mathrm{partition}}$', color='#fb8072', edgecolor='black')
    ax1.bar(x + width, p_net, width, label=r'$P_{\mathrm{net\ flux}}$', color='#80b1d3', edgecolor='black')
    ax1.set_ylabel('Relative probability', fontsize=10)
    ax1.set_title('A) Individual component probabilities', fontsize=11)
    ax1.set_xticks(x)
    ax1.set_xticklabels([f'{c}\n({"BBB+" if c in ("Caffeine","Diazepam","Morphine","Heroin") else "BBB-"})' for c in cases], fontsize=8)
    ax1.legend(fontsize=8)
    ax1.set_ylim(0, 1.05)

    colors = ['#1a9850' if c in ('Caffeine', 'Diazepam', 'Morphine', 'Heroin') else '#d73027' for c in cases]
    bars = ax2.bar(x, p_bbb, color=colors, edgecolor='black')
    ax2.set_ylabel(r'$P_{BBB}$ (product)', fontsize=10)
    ax2.set_title(r'B) $P_{BBB} = P_{\mathrm{desolv}} \times P_{\mathrm{partition}} \times P_{\mathrm{net\ flux}}$', fontsize=11)
    ax2.set_xticks(x)
    ax2.set_xticklabels(cases, rotation=45, ha='right', fontsize=9)
    ax2.axhline(0.1, color='red', linestyle='--', linewidth=1, alpha=0.5, label='Illustrative threshold')
    ax2.legend(fontsize=8)

    # label bars with values
    for bar, val in zip(bars, p_bbb):
        ax2.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.02,
                 f'{val:.2f}', ha='center', va='bottom', fontsize=8)

    fig.suptitle('Figure 5. Unified model decomposition for six clinical cases', fontsize=12, fontweight='bold', y=1.02)
    fig.tight_layout()
    path = output_dir / 'figure5_clinical_paradoxes.png'
    fig.savefig(path, dpi=get_param(params, 'FIGURE_DPI', as_float=True), bbox_inches='tight')
    plt.close(fig)
    return path


# ---------------------------------------------------------------------------
# Figure 6: Applications framework
# ---------------------------------------------------------------------------
def make_figure6(output_dir, params):
    fig, ax = plt.subplots(figsize=(10, 5.5))
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 6)
    ax.axis('off')

    def box(x, y, w, h, title, body, color, title_color='white', body_color='white', title_size=12, body_size=10):
        rect = FancyBboxPatch((x, y), w, h, boxstyle='round,pad=0.06',
                              facecolor=color, edgecolor='#333333', linewidth=2)
        ax.add_patch(rect)
        ax.text(x + w/2, y + h - 0.25, title, ha='center', va='top',
                fontsize=title_size, fontweight='bold', color=title_color)
        # body text
        lines = body.split('\n')
        start_y = y + h - 0.7
        for line in lines:
            ax.text(x + w/2, start_y, line, ha='center', va='top', fontsize=body_size, color=body_color)
            start_y -= 0.35

    box(0.5, 3.3, 2.6, 2.2, 'Micellar formulation',
        'Rate-limiting step becomes\ncarrier release / uptake\nand targeting ligands',
        '#45b7d1')
    box(3.7, 3.3, 2.6, 2.2, 'Metal-chelator delivery',
        'CNS delivery depends on\nformulation and route, not\nonly intrinsic descriptors',
        '#f39c12')
    box(6.9, 3.3, 2.6, 2.2, 'Local-anesthetic design',
        'Favor high $A_D$, shield HBDs\nbut low penetration $\\neq$ no\nCNS effect (affinity matters)',
        '#58b19d')

    # arrows to unified model
    for x in [1.8, 5.0, 8.2]:
        ax.annotate('', xy=(5.0, 2.3), xytext=(x, 3.3),
                    arrowprops=dict(arrowstyle='->', lw=2, color='#555555'))

    box(1.5, 0.7, 7.0, 1.6, 'Indicative, hypothesis-generating applications',
        r'$P_{BBB} \propto P_{\mathrm{desolv}} \times P_{\mathrm{partition}} \times P_{\mathrm{net\ flux}}$',
        '#2c3e50', title_color='#f1c40f', body_color='#f1c40f', title_size=13, body_size=11)

    ax.text(5.0, 5.7, 'Figure 6. Conceptual framework for applying the unified model to drug-design strategies',
            ha='center', fontsize=13, fontweight='bold')

    fig.tight_layout()
    path = output_dir / 'figure6_applications.png'
    fig.savefig(path, dpi=get_param(params, 'FIGURE_DPI', as_float=True), bbox_inches='tight')
    plt.close(fig)
    return path


# ---------------------------------------------------------------------------
# Top-level generator
# ---------------------------------------------------------------------------
def make_table3(df, params):
    """Return a Markdown table explaining six clinical cases within the unified model."""
    scores = _clinical_case_scores(df, params)
    rows = [
        ['Phenomenon', 'Conventional expectation', 'Unified model explanation'],
        [
            'Caffeine BBB+ (logP = -0.07)',
            'Difficult to explain by lipophilicity alone',
            'Small MW and $A_D$ allow paracellular/small-molecule diffusion; HBD = 0 gives high $P_\\mathrm{desolv}$ and $P_\\mathrm{net\\ flux}$.'
        ],
        [
            'Loperamide BBB- (logP = 4.77)',
            'High logP should favour entry',
            'Very large $A_D$ (~90 Å²) lowers $P_\\mathrm{partition}$; strong P-gp efflux makes $P_\\mathrm{net\\ flux}$ very low.'
        ],
        [
            'Morphine low BBB (HBD = 2)',
            'Moderate logP suggests moderate entry',
            'Two hydroxyl HBDs raise $P_\\mathrm{desolv}$ cost; it is also a P-gp substrate, lowering $P_\\mathrm{net\\ flux}$.'
        ],
        [
            'Heroin BBB+ vs morphine',
            'Higher MW should reduce entry',
            'Acetylation removes HBDs and lowers desolvation cost; $P_\\mathrm{desolv}$ becomes much larger, despite similar $A_D$.'
        ],
        [
            'Diazepam BBB+ vs loratadine BBB-',
            'Similar logP/MW but opposite BBB status',
            'Diazepam is not a P-gp substrate ($P_\\mathrm{net\\ flux}$ high); loratadine is a P-gp substrate and has larger $A_D$.'
        ],
        [
            'Dopamine BBB- (small, polar)',
            'Small size should aid entry',
            'Charged catecholamine at physiological pH gives very high $P_\\mathrm{desolv}$; no LAT1-like transporter dominates in this dataset.'
        ],
    ]
    return _markdown_table(rows, [26, 34, 70])


def compute_clinical_case_scores(df, params):
    """Public alias for _clinical_case_scores; used by build.py for Table 3 and other outputs."""
    return _clinical_case_scores(df, params)


def generate_figures(df, params, table2_rows, output_dir):
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    fig1 = make_figure1(output_dir, table2_rows, params)
    fig2 = make_figure2(output_dir, df, params)
    fig3 = make_figure3(output_dir, df, params)
    fig4 = make_figure4(output_dir, params)
    fig5 = make_figure5(output_dir, df, params)
    fig6 = make_figure6(output_dir, params)

    fig_paths = [fig1, fig2, fig3, fig4, fig5, fig6]
    titles = [
        'Figure 1. Discriminatory power ranking of unconventional and conventional descriptors',
        'Figure 2. Drug × factor evaluation matrix',
        'Figure 3. BBB permeability as a function of estimated A_D',
        'Figure 4. Unified three-component model of BBB permeability',
        'Figure 5. Unified model decomposition for six clinical cases',
        'Figure 6. Conceptual framework for applying the unified model to drug-design strategies',
    ]
    captions = [
        'Rating derived from the discriminatory-power assessment in Table 2; values are schematic overall promise ratings.',
        'Symbols: + favorable, o neutral, − weak, x unfavorable.  Rows are ordered by BBB status.',
        f'Estimated A_D values for the 24-drug dataset. Dashed lines mark the A_D ≈ {int(get_param(params, "AD_LOW_A2", as_float=True))}, {int(get_param(params, "AD_CUTOFF_A2", as_float=True))}, and {int(get_param(params, "AD_HIGH_A2", as_float=True))} Å² boundaries.',
        'A molecule must pass desolvation, membrane partition, and net transmembrane-flux gates.',
        'Illustrative relative probabilities derived from the heuristic scoring rules described in Methods.',
        'Applications to micellar encapsulation, metal-chelator delivery, and local-anesthetic design are indicative and require experimental validation.',
    ]
    pptx_path = output_dir / 'figures.pptx'
    _save_pptx(fig_paths, titles, captions, pptx_path, params)

    return {
        'fig1': fig1,
        'fig2': fig2,
        'fig3': fig3,
        'fig4': fig4,
        'fig5': fig5,
        'fig6': fig6,
        'pptx': pptx_path,
    }
