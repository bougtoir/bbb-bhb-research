"""Generate a revision-aligned graphical abstract for the EJPS BBB manuscript.

Outputs:
- output/graphical_abstract.png (high-resolution raster, suitable for journal upload)
- output/graphical_abstract.pptx (editable Python-pptx wrapper with title and caption)

The script is self-contained and can be rerun as part of the build pipeline.
"""
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


def make_graphical_abstract(output_dir):
    """Create the graphical abstract PNG and PPTX."""
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    png_path = _make_png(output_dir)
    pptx_path = _make_pptx(output_dir, png_path)
    return png_path, pptx_path


def _make_png(output_dir):
    fig_w, fig_h = 13.333, 7.5
    dpi = 300
    fig, ax = plt.subplots(figsize=(fig_w, fig_h))
    ax.set_xlim(0, 100)
    ax.set_ylim(0, 100)
    ax.axis('off')
    fig.patch.set_facecolor('white')
    ax.set_facecolor('white')

    def add_box(x, y, w, h, title, body, facecolor, title_color='white',
                body_color='white', title_fs=10, body_fs=9):
        rect = FancyBboxPatch((x, y), w, h, boxstyle='round,pad=0.5,rounding_size=1.0',
                              facecolor=facecolor, edgecolor='#2B2D42', linewidth=2)
        ax.add_patch(rect)
        ax.text(x + w/2, y + h - 3.5, title, ha='center', va='top',
                fontsize=title_fs, color=title_color, fontweight='bold')
        body_text = '\n'.join(body)
        ax.text(x + w/2, y + h/2 - 2, body_text, ha='center', va='center',
                fontsize=body_fs, color=body_color, linespacing=1.35)

    # Header
    header = FancyBboxPatch((0, 92), 100, 8, boxstyle='square,pad=0',
                            facecolor='#2B2D42', edgecolor='none')
    ax.add_patch(header)
    ax.text(50, 96, 'A Unified Three-Gate Model of BBB Permeability',
            ha='center', va='center', fontsize=22, color='white', fontweight='bold')

    # Top row
    y_top, h_top = 40, 45
    x_cases, w_cases = 3, 16
    add_box(x_cases, y_top, w_cases, h_top, 'Beyond simple rules',
            ['Caffeine: BBB+', 'paracellular diffusion',
             'Loperamide: BBB-', 'strong P-gp efflux',
             'Morphine vs Heroin', 'desolvation cost'],
            facecolor='#F28482')

    w_gate = 13
    x_g1, x_g2, x_g3 = 21, 36, 51
    add_box(x_g1, y_top, w_gate, h_top, 'Step 1: Desolvation',
            ['HBD / HBA', 'ΔSolv cost', '3D-PSA / charge'],
            facecolor='#3DCCC7', title_fs=9.5)
    add_box(x_g2, y_top, w_gate, h_top, 'Step 2: Partition',
            ['Membrane A_D', 'lateral pressure π_bi', 'logP'],
            facecolor='#4EA5DE', title_fs=9.5)
    add_box(x_g3, y_top, w_gate, h_top, 'Step 3: Net flux',
            ['P-gp efflux', 'Influx − Efflux', 'Transporters'],
            facecolor='#9B5DE5', title_fs=9.5)

    arrow_y = 62.5
    ax.annotate('', xy=(x_g1, arrow_y), xytext=(x_cases + w_cases, arrow_y),
                arrowprops=dict(arrowstyle='->', color='#2B2D42', lw=2))
    ax.annotate('', xy=(x_g2, arrow_y), xytext=(x_g1 + w_gate, arrow_y),
                arrowprops=dict(arrowstyle='->', color='#2B2D42', lw=2))
    ax.annotate('', xy=(x_g3, arrow_y), xytext=(x_g2 + w_gate, arrow_y),
                arrowprops=dict(arrowstyle='->', color='#2B2D42', lw=2))

    # Multiplication symbols are implicit in the sequential gate model; no extra markers are drawn on the arrows.

    x_unified, w_unified = 67, 13
    ax.annotate('', xy=(x_unified, arrow_y), xytext=(x_g3 + w_gate, arrow_y),
                arrowprops=dict(arrowstyle='->', color='#2B2D42', lw=2))
    add_box(x_unified, 45, w_unified, 35, 'Unified Model',
            [r'$P_{BBB} \propto P_{desolv}$',
             r'$\times P_{partition}$',
             r'$\times P_{net\ flux}$'],
            facecolor='#2B2D42', title_color='#FFD166', title_fs=11, body_fs=9)

    # Application label
    ax.text(50, 36, 'Indicative, hypothesis-generating applications',
            ha='center', va='center', fontsize=12, color='#2B2D42', style='italic')

    # Application boxes
    app_w, app_h, app_y, gap = 28, 24, 7, 2.5
    colors = ['#4CC9F0', '#F9A620', '#70C1B3']
    titles = ['Micellar formulation', 'Metal-chelator delivery', 'Local-anesthetic design']
    texts = [
        ['Rate-limiting step becomes', 'carrier release / uptake', 'and targeting ligands'],
        ['CNS delivery depends on', 'formulation and route, not', 'only intrinsic descriptors'],
        ['Favor high A_D, shield HBDs', 'but low penetration ≠ no', 'CNS effect (affinity matters)']
    ]
    app_centers = []
    for i, (title, body, color) in enumerate(zip(titles, texts, colors)):
        x = 3 + i * (app_w + gap)
        add_box(x, app_y, app_w, app_h, title, body, facecolor=color, body_fs=9.5)
        app_centers.append(x + app_w/2)

    # Arrows from gates to applications
    for gx, ac in zip([x_g1 + w_gate/2, x_g2 + w_gate/2, x_g3 + w_gate/2], app_centers):
        ax.annotate('', xy=(ac, app_y + app_h), xytext=(gx, y_top),
                    arrowprops=dict(arrowstyle='->', color='#2B2D42', lw=1.5, alpha=0.5))

    png_path = output_dir / 'graphical_abstract.png'
    fig.savefig(png_path, dpi=dpi, bbox_inches='tight', pad_inches=0.1)
    plt.close(fig)
    return png_path


def _make_pptx(output_dir, png_path):
    """Create an editable PPTX wrapper (title + image + caption)."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.6))
    tf = title_box.text_frame
    tf.text = 'Graphical Abstract'
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Caption
    caption_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12.333), Inches(0.5))
    tf2 = caption_box.text_frame
    tf2.text = ('A unified three-gate model of BBB permeability with illustrative clinical examples '
                'and hypothesis-generating applications.')
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.font.size = Pt(12)
    p2.alignment = PP_ALIGN.CENTER

    # Image centered in remaining area; preserve aspect ratio by setting height only
    img_h = Inches(5.3)
    top = Inches(1.0)
    pic = slide.shapes.add_picture(str(png_path), left=0, top=top, height=img_h)
    pic.left = int((prs.slide_width - pic.width) / 2)

    pptx_path = output_dir / 'graphical_abstract.pptx'
    prs.save(str(pptx_path))
    return pptx_path


def main():
    import sys
    out = Path(sys.argv[1]) if len(sys.argv) > 1 else Path(__file__).resolve().parent.parent / 'output'
    png, pptx = make_graphical_abstract(out)
    print(f'Graphical abstract written to:\n  {png}\n  {pptx}')


if __name__ == '__main__':
    main()
